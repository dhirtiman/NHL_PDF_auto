import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import os

# File paths
pdf_path = "./pdf_ppt/input/modules.pdf"
image_folder = "./pdf_ppt/images"
pptx_path = "./pdf_ppt/output/KineSpark_understandingGUTHealth.pptx"

# Create folder for images
os.makedirs(image_folder, exist_ok=True)

# Load PDF and convert pages to images
doc = fitz.open(pdf_path)
image_paths = []

print("Converting PDF pages to images...")
for i, page in enumerate(doc):
    pix = page.get_pixmap(dpi=100)
    img_path = f"{image_folder}/page_{i+1}.jpg"
    pix.save(img_path, "jpeg")
    image_paths.append(img_path)





# Create PowerPoint presentation
prs = Presentation()

# Set slide dimensions to 16:9 widescreen
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

print("Creating PPT slides...")
for img_path in image_paths:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)

prs.save(pptx_path)
print(f"✅ PPT created successfully: {pptx_path}")
